# ruff: noqa: E501
"""Generate document bytes and upload them to object storage."""

from __future__ import annotations

import os
import re
import textwrap
from datetime import UTC, datetime
from html import escape as html_escape
from io import BytesIO
from typing import Any

import httpx
from temporalio import activity

from worker.activities.document_generation.types import (
    MIME_TYPES,
    DocumentGenerationDestination,
    DocumentGenerationRequest,
    DocumentGenerationSourceBundle,
    DocumentGenerationSourceItem,
    DocumentGenerationWorkflowParams,
    GeneratedDocumentArtifact,
    normalize_source_bundle,
)
from worker.lib.s3_client import upload_bytes

PDF_PAGE_WIDTH = 595
PDF_PAGE_HEIGHT = 842
PDF_MARGIN_LEFT = 54
PDF_MARGIN_TOP = 56
PDF_MARGIN_BOTTOM = 56
PDF_FONT_SIZE = 11
PDF_LINE_HEIGHT = 16
PDF_CJK_FONT = "korea"
MAX_SOURCE_CHARS = 4000
PPTX_WRAP_CHARS = 82
PPTX_LINES_PER_SLIDE = 8
TECTONIC_URL = os.environ.get("TECTONIC_URL", "http://tectonic:8888")
TECTONIC_TIMEOUT_S = float(os.environ.get("TECTONIC_TIMEOUT_S", "120"))
TECTONIC_ERROR_SUMMARY_CHARS = 2000
_tectonic_client: httpx.AsyncClient | None = None

TEMPLATE_LABELS: dict[str, dict[str, str]] = {
    "report": {
        "prompt": "Objective",
        "context": "Generation Context",
        "sources": "Source Register",
        "quality": "Source Quality Notes",
        "body": "Source Detail",
    },
    "brief": {
        "prompt": "Brief",
        "context": "Generation Context",
        "sources": "Source Register",
        "quality": "Source Quality Notes",
        "body": "Supporting Notes",
    },
    "research_summary": {
        "prompt": "Research Question",
        "context": "Generation Context",
        "sources": "Evidence Register",
        "quality": "Evidence Quality Notes",
        "body": "Evidence Detail",
    },
    "deck": {
        "prompt": "Deck Goal",
        "context": "Generation Context",
        "sources": "Slide Source Register",
        "quality": "Source Quality Notes",
        "body": "Slide Notes",
    },
    "spreadsheet": {
        "prompt": "Workbook Goal",
        "context": "Generation Context",
        "sources": "Data Source Register",
        "quality": "Data Quality Notes",
        "body": "Source Detail",
    },
    "custom": {
        "prompt": "Custom Instructions",
        "context": "Generation Context",
        "sources": "Source Register",
        "quality": "Source Quality Notes",
        "body": "Source Detail",
    },
    "technical_report": {
        "prompt": "Objective",
        "context": "Report Context",
        "sources": "Source Register",
        "quality": "Source Quality Notes",
        "body": "Technical Detail",
    },
    "research_brief": {
        "prompt": "Research Question",
        "context": "Brief Context",
        "sources": "Evidence Register",
        "quality": "Evidence Quality Notes",
        "body": "Evidence Detail",
    },
    "paper_style": {
        "prompt": "Abstract",
        "context": "Paper Context",
        "sources": "Evidence Register",
        "quality": "Evidence Quality Notes",
        "body": "Findings",
    },
    "business_report": {
        "prompt": "Executive Summary",
        "context": "Business Context",
        "sources": "Input Register",
        "quality": "Input Quality Notes",
        "body": "Analysis",
    },
}


def _value(data: dict[str, Any], snake: str, camel: str | None = None, default: Any = None) -> Any:
    if snake in data:
        return data[snake]
    if camel and camel in data:
        return data[camel]
    return default


def normalize_destination(
    raw: DocumentGenerationDestination | dict[str, Any],
) -> DocumentGenerationDestination:
    if isinstance(raw, DocumentGenerationDestination):
        return raw
    return DocumentGenerationDestination(
        filename=str(raw["filename"]),
        title=raw.get("title"),
        folder_id=_value(raw, "folder_id", "folderId"),
        publish_as=_value(raw, "publish_as", "publishAs", "agent_file"),
        start_ingest=bool(_value(raw, "start_ingest", "startIngest", False)),
    )


def normalize_generation(
    raw: DocumentGenerationRequest | dict[str, Any],
) -> DocumentGenerationRequest:
    if isinstance(raw, DocumentGenerationRequest):
        return raw
    return DocumentGenerationRequest(
        format=raw["format"],
        prompt=raw["prompt"],
        render_engine=_value(raw, "render_engine", "renderEngine"),
        image_engine=_value(raw, "image_engine", "imageEngine"),
        locale=raw.get("locale", "ko"),
        template=raw.get("template", "report"),
        sources=list(raw.get("sources", [])),
        destination=normalize_destination(raw["destination"]),
        artifact_mode=_value(raw, "artifact_mode", "artifactMode", "object_storage"),
    )


def normalize_params(
    raw: DocumentGenerationWorkflowParams | dict[str, Any],
) -> DocumentGenerationWorkflowParams:
    if isinstance(raw, DocumentGenerationWorkflowParams):
        generation = normalize_generation(raw.generation)
        return DocumentGenerationWorkflowParams(
            action_id=raw.action_id,
            request_id=raw.request_id,
            workspace_id=raw.workspace_id,
            project_id=raw.project_id,
            user_id=raw.user_id,
            generation=generation,
        )
    return DocumentGenerationWorkflowParams(
        action_id=_value(raw, "action_id", "actionId"),
        request_id=_value(raw, "request_id", "requestId"),
        workspace_id=_value(raw, "workspace_id", "workspaceId"),
        project_id=_value(raw, "project_id", "projectId"),
        user_id=_value(raw, "user_id", "userId"),
        generation=normalize_generation(raw["generation"]),
    )


def _safe_path_segment(value: str) -> str:
    normalized = re.sub(r"[^A-Za-z0-9._-]+", "-", value.strip()).strip(".-")
    return normalized[:120] or "document"


def _artifact_key(
    params: DocumentGenerationWorkflowParams, generation: DocumentGenerationRequest
) -> str:
    destination = normalize_destination(generation.destination)
    filename = _safe_path_segment(destination.filename)
    request_id = _safe_path_segment(params.request_id)
    project_id = _safe_path_segment(params.project_id)
    return f"agent-files/{project_id}/document-generation/{request_id}/{filename}"


def _source_reference_lines(generation: DocumentGenerationRequest) -> list[str]:
    lines = []
    for source in generation.sources:
        label = source.get("type", "source")
        source_id = (
            source.get("noteId")
            or source.get("objectId")
            or source.get("threadId")
            or source.get("runId")
            or "unknown"
        )
        lines.append(f"- {label}: {source_id}")
    return lines


def _truncate_cell_text(value: str, limit: int = MAX_SOURCE_CHARS) -> str:
    normalized = value.strip()
    if len(normalized) <= limit:
        return normalized
    return f"{normalized[:limit].rstrip()}\n..."


def _template_labels(template: str) -> dict[str, str]:
    return TEMPLATE_LABELS.get(template, TEMPLATE_LABELS["custom"])


def _source_summary(source: DocumentGenerationSourceItem) -> str:
    status = "included" if source.included else "omitted"
    signal_text = (
        f"; signals: {', '.join(dict.fromkeys(source.quality_signals))}"
        if source.quality_signals
        else ""
    )
    return f"{source.kind}: {source.title} ({source.id}; {status}{signal_text})"


def _source_quality_bullets(sources: list[DocumentGenerationSourceItem]) -> list[str]:
    bullets: list[str] = []
    for source in sources:
        signals = list(dict.fromkeys(source.quality_signals))
        if not source.included:
            signals.append("source_omitted")
        if signals:
            bullets.append(f"{source.kind}: {source.title} - {', '.join(dict.fromkeys(signals))}")
    return bullets


def _source_excerpt(source: DocumentGenerationSourceItem, *, chars: int) -> str:
    if not source.included:
        return "Source omitted from body because it exceeded the document generation budget."
    return _truncate_cell_text(source.body, chars)


def _document_model(
    params: DocumentGenerationWorkflowParams,
    generation: DocumentGenerationRequest,
    sources: DocumentGenerationSourceBundle,
) -> dict[str, Any]:
    destination = normalize_destination(generation.destination)
    title = destination.title or destination.filename
    included_sources = [source for source in sources.items if source.included]
    all_sources = list(sources.items)
    reference_lines = _source_reference_lines(generation)
    labels = _template_labels(generation.template)
    context_bullets = [
        f"Format: {generation.format}",
        f"Template: {generation.template}",
        f"Locale: {generation.locale}",
        f"Request: {params.request_id}",
    ]
    source_quality_bullets = _source_quality_bullets(all_sources)
    sections = [
        {
            "title": labels["prompt"],
            "body": generation.prompt,
            "bullets": [],
            "kind": "prompt",
        },
        {
            "title": labels["context"],
            "body": "",
            "bullets": context_bullets,
            "kind": "context",
        },
    ]
    if all_sources:
        sections.append(
            {
                "title": labels["sources"],
                "body": "",
                "bullets": [_source_summary(source) for source in all_sources],
                "kind": "sources",
            }
        )
        if source_quality_bullets:
            sections.append(
                {
                    "title": labels["quality"],
                    "body": "",
                    "bullets": source_quality_bullets,
                    "kind": "quality",
                }
            )
        sections.extend(
            {
                "title": f"{labels['body']}: {source.title}",
                "body": _source_excerpt(
                    source,
                    chars=1800 if generation.template == "brief" else MAX_SOURCE_CHARS,
                ),
                "bullets": [],
                "kind": "source_body",
            }
            for source in included_sources
        )
    else:
        sections.append(
            {
                "title": labels["sources"],
                "body": "",
                "bullets": reference_lines or ["No explicit sources provided"],
                "kind": "sources",
            }
        )
    return {
        "title": title,
        "prompt": generation.prompt,
        "format": generation.format,
        "template": generation.template,
        "locale": generation.locale,
        "request_id": params.request_id,
        "sections": sections,
        "sources": included_sources,
        "all_sources": all_sources,
        "quality_bullets": source_quality_bullets,
    }


def _render_pdf_model(model: dict[str, Any]) -> bytes:
    import pymupdf

    doc = pymupdf.open()
    page = None
    y = PDF_MARGIN_TOP
    page_bottom = PDF_PAGE_HEIGHT - PDF_MARGIN_BOTTOM
    font = pymupdf.Font(PDF_CJK_FONT)

    def new_page() -> Any:
        nonlocal page, y
        page = doc.new_page(width=PDF_PAGE_WIDTH, height=PDF_PAGE_HEIGHT)
        y = PDF_MARGIN_TOP
        return page

    def overflow_text(overflow: list[tuple[str, float]]) -> str:
        return " ".join(token for token, _width in overflow).strip()

    def split_oversized_token(text: str, *, size: int, width: float) -> tuple[str, str]:
        head = ""
        tail_start = 0
        for index, char in enumerate(text):
            candidate = f"{head}{char}"
            if head and font.text_length(candidate, fontsize=size) > width:
                tail_start = index
                break
            head = candidate
        else:
            return text, ""
        return head, text[tail_start:]

    def add_line(
        line: str,
        *,
        size: int = PDF_FONT_SIZE,
        line_height: int = PDF_LINE_HEIGHT,
        indent: int = 0,
        after: int = 0,
    ) -> None:
        nonlocal page, y
        available_width = PDF_PAGE_WIDTH - PDF_MARGIN_LEFT - PDF_MARGIN_LEFT - indent
        remaining = line
        if not remaining:
            if page is None or y + line_height > page_bottom:
                new_page()
            y += line_height
            y += after
            return
        while remaining:
            if page is None or y + line_height > page_bottom:
                new_page()
            rect = pymupdf.Rect(
                PDF_MARGIN_LEFT + indent,
                y,
                PDF_PAGE_WIDTH - PDF_MARGIN_LEFT,
                page_bottom,
            )
            writer = pymupdf.TextWriter(page.rect)
            overflow = writer.fill_textbox(
                rect,
                remaining,
                font=font,
                fontsize=size,
                lineheight=line_height,
            )
            wrote_text = not overflow or writer.text_rect.height > 0
            if wrote_text:
                writer.write_text(page)
                y = max(y + line_height, min(writer.text_rect.y1 + 2, page_bottom + 1))
                remaining = overflow_text(overflow)
                if remaining:
                    new_page()
                continue

            head, tail = split_oversized_token(remaining, size=size, width=available_width)
            if not head:
                raise ValueError("document_generation_pdf_text_does_not_fit")
            writer = pymupdf.TextWriter(page.rect)
            overflow = writer.fill_textbox(
                rect,
                head,
                font=font,
                fontsize=size,
                lineheight=line_height,
            )
            if overflow:
                raise ValueError("document_generation_pdf_text_does_not_fit")
            writer.write_text(page)
            y = max(y + line_height, min(writer.text_rect.y1 + 2, page_bottom + 1))
            remaining = tail
            if remaining:
                new_page()
        y += after

    add_line(str(model["title"]), size=16, line_height=22, after=10)
    add_line(
        f"Template: {model['template']} / Locale: {model['locale']} / Request: {model['request_id']}",
        size=9,
        line_height=13,
        after=12,
    )
    for section in model["sections"]:
        add_line(str(section["title"]), size=13, line_height=19, after=3)
        if section.get("body"):
            for paragraph in str(section["body"]).splitlines() or [""]:
                add_line(paragraph, size=PDF_FONT_SIZE, line_height=PDF_LINE_HEIGHT)
            y += 5
        for bullet in section.get("bullets", []):
            add_line(f"- {bullet}", size=PDF_FONT_SIZE, line_height=PDF_LINE_HEIGHT, indent=10)
        y += 7

    return doc.tobytes(garbage=4, deflate=True)


def _render_docx(model: dict[str, Any]) -> bytes:
    from docx import Document
    from docx.enum.text import WD_BREAK
    from docx.shared import Pt

    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Malgun Gothic"
    styles["Normal"].font.size = Pt(10.5)
    document.add_heading(str(model["title"]), level=0)
    meta = document.add_table(rows=0, cols=2)
    for field, value in (
        ("Template", model["template"]),
        ("Locale", model["locale"]),
        ("Request", model["request_id"]),
    ):
        cells = meta.add_row().cells
        cells[0].text = field
        cells[1].text = str(value)
    for index, section in enumerate(model["sections"]):
        document.add_heading(str(section["title"]), level=1)
        if section.get("body"):
            paragraphs = str(section["body"]).splitlines() or [""]
            for paragraph in paragraphs:
                document.add_paragraph(paragraph)
        for bullet in section.get("bullets", []):
            document.add_paragraph(str(bullet), style="List Bullet")
        if section.get("kind") == "source_body" and index < len(model["sections"]) - 1:
            document.paragraphs[-1].runs[-1].add_break(WD_BREAK.PAGE)

    out = BytesIO()
    document.save(out)
    return out.getvalue()


def _chunked(values: list[str], size: int) -> list[list[str]]:
    if not values:
        return [[]]
    return [values[index : index + size] for index in range(0, len(values), size)]


def _pptx_wrapped_lines(values: list[str]) -> list[str]:
    lines: list[str] = []
    for value in values:
        text = value.strip()
        if not text:
            lines.append("")
            continue
        for raw_line in text.splitlines():
            wrapped = textwrap.wrap(
                raw_line,
                width=PPTX_WRAP_CHARS,
                break_long_words=True,
                replace_whitespace=False,
                drop_whitespace=False,
            )
            lines.extend(wrapped or [""])
    return lines


def _render_pptx(model: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = str(model["title"])
    title_slide.placeholders[1].text = (
        f"{model['template']} / {model['locale']} / {len(model['sources'])} sources"
    )

    for section in model["sections"]:
        content = []
        if section.get("body"):
            content.append(str(section["body"]))
        content.extend(str(bullet) for bullet in section.get("bullets", []))
        for page_index, lines in enumerate(
            _chunked(_pptx_wrapped_lines(content), PPTX_LINES_PER_SLIDE)
        ):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title_suffix = f" ({page_index + 1})" if page_index else ""
            slide.shapes.title.text = f"{section['title']}{title_suffix}"
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, line in enumerate(lines):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = line
                paragraph.font.size = Pt(17 if section.get("kind") != "source_body" else 14)
                paragraph.level = 0

    if model.get("quality_bullets"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Quality Notes"
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(4.6))
        frame = box.text_frame
        frame.clear()
        for index, line in enumerate(_pptx_wrapped_lines(model["quality_bullets"])):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(14)

    out = BytesIO()
    presentation.save(out)
    return out.getvalue()


def _render_xlsx(model: dict[str, Any]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font

    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.append(["Field", "Value"])
    summary.append(["Title", model["title"]])
    summary.append(["Prompt", model["prompt"]])
    summary.append(["Template", model["template"]])
    summary.append(["Locale", model["locale"]])
    summary.append(["Request", model["request_id"]])
    summary.append(["Source Count", len(model["all_sources"])])
    summary.append(["Generated At", datetime.now(UTC).isoformat()])

    outline = workbook.create_sheet("Outline")
    outline.append(["Section", "Type", "Content"])
    for section in model["sections"]:
        if section.get("body"):
            outline.append(
                [section["title"], section.get("kind", "section"), _truncate_cell_text(section["body"])]
            )
        for bullet in section.get("bullets", []):
            outline.append([section["title"], "bullet", str(bullet)])

    sources = workbook.create_sheet("Sources")
    sources.append(["Source ID", "Kind", "Title", "Included", "Quality Signals", "Excerpt"])
    for source in model["all_sources"]:
        sources.append(
            [
                source.id,
                source.kind,
                source.title,
                source.included,
                ", ".join(dict.fromkeys(source.quality_signals)),
                _source_excerpt(source, chars=1000),
            ]
        )

    if model.get("quality_bullets"):
        quality = workbook.create_sheet("Quality")
        quality.append(["Source Quality Signal"])
        for bullet in model["quality_bullets"]:
            quality.append([bullet])

    for sheet in workbook.worksheets:
        sheet.freeze_panes = "A2"
        sheet.auto_filter.ref = sheet.dimensions
        for cell in sheet[1]:
            cell.font = Font(bold=True)
        for column in sheet.columns:
            letter = column[0].column_letter
            sheet.column_dimensions[letter].width = min(
                max(len(str(cell.value or "")) for cell in column) + 2,
                72,
            )
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

    out = BytesIO()
    workbook.save(out)
    return out.getvalue()


_LATEX_ESCAPE_RE = re.compile(r"([&%$#_{}])")


def _latex_escape(value: Any) -> str:
    text = str(value or "")
    text = text.replace("\\", r"\textbackslash{}")
    text = _LATEX_ESCAPE_RE.sub(r"\\\1", text)
    text = text.replace("~", r"\textasciitilde{}")
    text = text.replace("^", r"\textasciicircum{}")
    return text


def _latex_paragraphs(value: Any) -> str:
    paragraphs = []
    for paragraph in str(value or "").splitlines():
        if paragraph.strip():
            paragraphs.append(_latex_escape(paragraph))
    return "\n\n".join(paragraphs)


def _latex_template_preamble(template: str) -> str:
    base = r"""
\documentclass[a4paper,11pt]{article}
\usepackage{kotex}
\usepackage[a4paper,margin=1in]{geometry}
\usepackage{graphicx}
\usepackage{booktabs}
\usepackage{longtable}
\usepackage{hyperref}
\usepackage{enumitem}
\setlist[itemize]{leftmargin=1.5em}
"""
    if template == "paper_style":
        return base + r"\usepackage{abstract}" + "\n"
    return base


def _render_latex_source(model: dict[str, Any]) -> str:
    lines = [
        _latex_template_preamble(str(model["template"])),
        f"\\title{{{_latex_escape(model['title'])}}}",
        "\\date{\\today}",
        "\\begin{document}",
        "\\maketitle",
        "\\tableofcontents",
        "\\newpage",
    ]
    for section in model["sections"]:
        lines.append(f"\\section{{{_latex_escape(section['title'])}}}")
        if section.get("body"):
            lines.append(_latex_paragraphs(section["body"]))
        bullets = section.get("bullets", [])
        if bullets:
            lines.append("\\begin{itemize}")
            for bullet in bullets:
                lines.append(f"  \\item {_latex_escape(bullet)}")
            lines.append("\\end{itemize}")
    lines.append("\\end{document}")
    return "\n\n".join(line for line in lines if line)


async def _post_tectonic(tex_source: str, bib_source: str) -> bytes:
    global _tectonic_client
    if _tectonic_client is None or _tectonic_client.is_closed:
        _tectonic_client = httpx.AsyncClient(timeout=TECTONIC_TIMEOUT_S)
    response = await _tectonic_client.post(
        f"{TECTONIC_URL}/compile",
        json={
            "tex_source": tex_source,
            "bib_source": bib_source or None,
            "engine": "xelatex",
            "timeout_ms": int(TECTONIC_TIMEOUT_S * 1000),
        },
    )
    if response.status_code == 504:
        raise RuntimeError("tectonic_timeout")
    if response.status_code != 200:
        raise RuntimeError(
            f"tectonic_failed: {response.text[:TECTONIC_ERROR_SUMMARY_CHARS]}"
        )
    return response.content


async def _render_latex_pdf_model(model: dict[str, Any]) -> bytes:
    return await _post_tectonic(_render_latex_source(model), "")


def _svg_text_lines(text: str, *, width: int, max_lines: int) -> list[str]:
    lines: list[str] = []
    for raw_line in text.splitlines() or [text]:
        normalized = raw_line.strip()
        if not normalized:
            continue
        lines.extend(textwrap.wrap(normalized, width=width, break_long_words=True))
        if len(lines) >= max_lines:
            return [*lines[: max_lines - 1], f"{lines[max_lines - 1]}..."]
    return lines[:max_lines]


def _render_svg_text_block(
    *,
    title: str,
    lines: list[str],
    x: int,
    y: int,
    width: int,
    fill: str,
) -> str:
    safe_title = html_escape(title)
    text_lines = [
        f'<text x="{x}" y="{y}" font-size="28" font-weight="700" fill="{fill}">{safe_title}</text>'
    ]
    current_y = y + 38
    for line in lines:
        text_lines.append(
            f'<text x="{x}" y="{current_y}" font-size="22" fill="{fill}" opacity="0.9">{html_escape(line)}</text>'
        )
        current_y += 32
    return (
        f'<rect x="{x - 24}" y="{y - 38}" width="{width}" height="{current_y - y + 28}" rx="18" '
        'fill="rgba(255,255,255,0.88)" stroke="rgba(15,23,42,0.12)"/>'
        + "\n"
        + "\n".join(text_lines)
    )


def _render_svg_figure_model(model: dict[str, Any]) -> bytes:
    title = str(model["title"])
    prompt = str(model["prompt"])
    context = next(
        (section for section in model["sections"] if section.get("kind") == "context"),
        None,
    )
    sources = next(
        (section for section in model["sections"] if section.get("kind") == "sources"),
        None,
    )
    prompt_lines = _svg_text_lines(prompt, width=68, max_lines=5)
    context_lines = _svg_text_lines(
        "\n".join(str(item) for item in (context or {}).get("bullets", [])[:4]),
        width=48,
        max_lines=4,
    )
    source_lines = _svg_text_lines(
        "\n".join(str(item) for item in (sources or {}).get("bullets", [])[:4]),
        width=52,
        max_lines=5,
    )
    title_lines = _svg_text_lines(title, width=42, max_lines=2)
    safe_title = html_escape(title_lines[0] if title_lines else title)
    subtitle = html_escape(title_lines[1] if len(title_lines) > 1 else str(model["template"]))
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" viewBox="0 0 1600 900" role="img" aria-label="{safe_title}">
  <rect width="1600" height="900" fill="#f6f8fb"/>
  <rect x="0" y="0" width="1600" height="900" fill="#eaf1f4"/>
  <circle cx="1370" cy="120" r="220" fill="#d8e7df"/>
  <circle cx="180" cy="760" r="260" fill="#e2dfef"/>
  <text x="96" y="120" font-family="Inter, Arial, sans-serif" font-size="54" font-weight="800" fill="#0f172a">{safe_title}</text>
  <text x="100" y="166" font-family="Inter, Arial, sans-serif" font-size="24" fill="#475569">{subtitle}</text>
  <g font-family="Inter, Arial, sans-serif">
    {_render_svg_text_block(title="Objective", lines=prompt_lines, x=120, y=270, width=650, fill="#0f172a")}
    {_render_svg_text_block(title="Context", lines=context_lines, x=900, y=270, width=560, fill="#12323a")}
    {_render_svg_text_block(title="Sources", lines=source_lines, x=900, y=520, width=560, fill="#263247")}
  </g>
  <text x="100" y="820" font-family="Inter, Arial, sans-serif" font-size="18" fill="#64748b">OpenCairn generated figure · request {html_escape(str(model["request_id"]))}</text>
</svg>
"""
    return svg.encode("utf-8")


def _model_image_prompt(model: dict[str, Any]) -> str:
    source_summaries = [
        _source_summary(source)
        for source in model.get("all_sources", [])[:6]
    ]
    quality_notes = [
        str(item)
        for item in model.get("quality_bullets", [])[:4]
    ]
    parts = [
        "Create a polished, presentation-ready figure for an OpenCairn knowledge artifact.",
        "Use a clean professional visual style, balanced spacing, and legible labels.",
        "Avoid fictional citations. If evidence is uncertain, represent it as an uncertainty cue.",
        f"Title: {model['title']}",
        f"Template: {model['template']}",
        f"Objective: {model['prompt']}",
    ]
    if source_summaries:
        parts.append("Sources:\n" + "\n".join(f"- {line}" for line in source_summaries))
    if quality_notes:
        parts.append("Source quality notes:\n" + "\n".join(f"- {line}" for line in quality_notes))
    return "\n\n".join(parts)


async def _render_model_image(model: dict[str, Any]) -> tuple[bytes, str]:
    from llm.factory import get_provider

    provider = get_provider()
    result = await provider.generate_image(_model_image_prompt(model))
    if result is None:
        raise RuntimeError("image_generation_not_supported")
    return result.image_bytes, result.mime_type


def _render_document_model_bytes(
    model: dict[str, Any],
    format: str,
) -> bytes:
    if format == "pptx":
        return _render_pptx(model)
    if format == "xlsx":
        return _render_xlsx(model)
    if format == "pdf":
        return _render_pdf_model(model)
    if format == "docx":
        return _render_docx(model)
    if format == "image":
        return _render_svg_figure_model(model)
    raise ValueError(f"Unsupported document generation format: {format}")


def render_document_bytes(
    params: DocumentGenerationWorkflowParams,
    source_bundle: DocumentGenerationSourceBundle | dict[str, Any] | None = None,
) -> bytes:
    generation = normalize_generation(params.generation)
    sources = normalize_source_bundle(source_bundle)
    model = _document_model(params, generation, sources)
    return _render_document_model_bytes(model, generation.format)


def heartbeat_safe(message: str) -> None:
    try:
        activity.heartbeat(message)
    except RuntimeError as exc:
        if str(exc) != "Not in activity context":
            raise


@activity.defn(name="generate_document_artifact")
async def generate_document_artifact(
    params: DocumentGenerationWorkflowParams | dict[str, Any],
    source_bundle: DocumentGenerationSourceBundle | dict[str, Any] | None = None,
) -> GeneratedDocumentArtifact:
    normalized = normalize_params(params)
    generation = normalize_generation(normalized.generation)
    if generation.artifact_mode != "object_storage":
        raise ValueError("document_generation_requires_object_storage")

    heartbeat_safe(f"generating {generation.format}")
    sources = normalize_source_bundle(source_bundle)
    model = _document_model(normalized, generation, sources)
    if generation.format == "pdf" and generation.render_engine == "latex":
        heartbeat_safe("compiling latex pdf")
        body = await _render_latex_pdf_model(model)
        heartbeat_safe("compiled latex pdf")
        mime_type = MIME_TYPES[generation.format]
    elif generation.format == "image" and generation.image_engine == "model":
        body, mime_type = await _render_model_image(model)
    else:
        body = _render_document_model_bytes(model, generation.format)
        mime_type = MIME_TYPES[generation.format]
    object_key = _artifact_key(normalized, generation)
    upload_bytes(object_key, body, mime_type)
    return GeneratedDocumentArtifact(
        objectKey=object_key,
        mimeType=mime_type,
        bytes=len(body),
    )
