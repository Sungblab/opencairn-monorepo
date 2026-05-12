from __future__ import annotations

import asyncio
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from unittest.mock import patch

import pymupdf
import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from worker.activities.document_generation.generate import generate_document_artifact
from worker.activities.document_generation.register import register_document_generation_result
from worker.activities.document_generation.sources import hydrate_document_generation_sources
from worker.activities.document_generation.types import GeneratedDocumentArtifact


def _request(format_: str = "pdf") -> dict:
    return {
        "actionId": "00000000-0000-4000-8000-000000000011",
        "requestId": "00000000-0000-4000-8000-000000000020",
        "workspaceId": "00000000-0000-4000-8000-000000000001",
        "projectId": "00000000-0000-4000-8000-000000000003",
        "userId": "00000000-0000-4000-8000-000000000004",
        "generation": {
            "format": format_,
            "prompt": "Generate a polished project report.",
            "locale": "ko",
            "template": "report",
            "sources": [
                {
                    "type": "note",
                    "noteId": "00000000-0000-4000-8000-000000000021",
                }
            ],
            "destination": {
                "filename": f"project-report.{format_}",
                "title": "Project report",
                "publishAs": "agent_file",
                "startIngest": False,
            },
            "artifactMode": "object_storage",
        },
    }


@pytest.mark.asyncio
async def test_generate_document_artifact_uploads_pdf_to_object_storage() -> None:
    uploaded: dict[str, object] = {}

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        result = await generate_document_artifact(_request("pdf"))

    assert result.objectKey.endswith("/project-report.pdf")
    assert result.objectKey.startswith(
        "agent-files/00000000-0000-4000-8000-000000000003/document-generation/"
    )
    assert result.mimeType == "application/pdf"
    assert result.bytes == len(uploaded["data"])
    assert uploaded["content_type"] == "application/pdf"
    assert bytes(uploaded["data"]).startswith(b"%PDF-")


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_fetches_note_bodies() -> None:
    calls: list[tuple[str, dict]] = []

    async def fake_post(path: str, body: dict) -> dict:
        calls.append((path, body))
        return {
            "id": body["source_id"],
            "title": "시장 조사 노트",
            "body": "핵심 근거: 한국어 문서 생성 품질을 검증합니다.",
            "kind": "note",
        }

    with patch("worker.activities.document_generation.sources.post_internal", fake_post):
        result = await hydrate_document_generation_sources(_request("pdf"))

    assert calls == [
        (
            "/api/internal/synthesis-export/fetch-source",
            {
                "source_id": "00000000-0000-4000-8000-000000000021",
                "kind": "note",
                "workspace_id": "00000000-0000-4000-8000-000000000001",
            },
        )
    ]
    assert result.items[0].title == "시장 조사 노트"
    assert "한국어 문서 생성 품질" in result.items[0].body


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_fetches_notes_with_bounded_concurrency() -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "note",
            "noteId": f"00000000-0000-4000-8000-0000000000{index:02d}",
        }
        for index in range(21, 27)
    ]
    in_flight = 0
    max_in_flight = 0

    async def fake_post(_path: str, body: dict) -> dict:
        nonlocal in_flight, max_in_flight
        in_flight += 1
        max_in_flight = max(max_in_flight, in_flight)
        await asyncio.sleep(0.01)
        in_flight -= 1
        return {
            "id": body["source_id"],
            "title": f"노트 {body['source_id'][-2:]}",
            "body": "동시 hydration 검증 본문",
            "kind": "note",
        }

    with patch("worker.activities.document_generation.sources.post_internal", fake_post):
        result = await hydrate_document_generation_sources(request)

    assert len(result.items) == 6
    assert max_in_flight > 1


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_fetches_non_note_sources() -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "agent_file",
            "objectId": "00000000-0000-4000-8000-000000000031",
        },
        {
            "type": "chat_thread",
            "threadId": "00000000-0000-4000-8000-000000000032",
        },
        {
            "type": "research_run",
            "runId": "00000000-0000-4000-8000-000000000033",
        },
        {
            "type": "synthesis_run",
            "runId": "00000000-0000-4000-8000-000000000034",
        },
    ]
    calls: list[tuple[str, dict]] = []

    async def fake_post(path: str, body: dict) -> dict:
        calls.append((path, body))
        source = body["source"]
        source_id = source.get("objectId") or source.get("threadId") or source.get("runId")
        return {
            "id": source_id,
            "title": f"{source['type']} title",
            "body": f"{source['type']} hydrated content",
            "kind": source["type"],
            "tokenCount": 12,
        }

    with patch("worker.activities.document_generation.sources.post_internal", fake_post):
        result = await hydrate_document_generation_sources(request)

    assert [call[0] for call in calls] == [
        "/api/internal/document-generation/hydrate-source",
        "/api/internal/document-generation/hydrate-source",
        "/api/internal/document-generation/hydrate-source",
        "/api/internal/document-generation/hydrate-source",
    ]
    assert all(call[1]["workspaceId"] == request["workspaceId"] for call in calls)
    assert all(call[1]["projectId"] == request["projectId"] for call in calls)
    assert all(call[1]["userId"] == request["userId"] for call in calls)
    bodies = {item.kind: item.body for item in result.items}
    assert bodies["agent_file"] == "agent_file hydrated content"
    assert bodies["chat_thread"] == "chat_thread hydrated content"
    assert bodies["research_run"] == "research_run hydrated content"
    assert bodies["synthesis_run"] == "synthesis_run hydrated content"


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_falls_back_per_failed_source() -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "agent_file",
            "objectId": "00000000-0000-4000-8000-000000000031",
        },
        {
            "type": "chat_thread",
            "threadId": "00000000-0000-4000-8000-000000000032",
        },
    ]

    async def fake_post(_path: str, body: dict) -> dict:
        if body["source"]["type"] == "agent_file":
            raise RuntimeError("hydration API unavailable")
        return {
            "id": body["source"]["threadId"],
            "title": "Thread",
            "body": "hydrated thread body",
            "kind": "chat_thread",
            "tokenCount": 10,
        }

    with patch("worker.activities.document_generation.sources.post_internal", fake_post):
        result = await hydrate_document_generation_sources(request)

    by_kind = {item.kind: item for item in result.items}
    assert by_kind["agent_file"].title == "Project object"
    assert by_kind["agent_file"].body == "agent_file: 00000000-0000-4000-8000-000000000031"
    assert by_kind["chat_thread"].body == "hydrated thread body"


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_extracts_pdf_agent_file_text() -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "agent_file",
            "objectId": "00000000-0000-4000-8000-000000000031",
        }
    ]
    pdf = pymupdf.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Binary PDF source body")
    with NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(pdf.tobytes())
        tmp_path = Path(tmp.name)
    pdf.close()

    async def fake_post(_path: str, body: dict) -> dict:
        return {
            "id": body["source"]["objectId"],
            "title": "Uploaded PDF",
            "body": "uploaded.pdf (pdf, 2048 bytes)",
            "kind": "agent_file",
            "mimeType": "application/pdf",
            "objectKey": "agent-files/uploaded.pdf",
            "bytes": 2048,
        }

    with (
        patch("worker.activities.document_generation.sources.post_internal", fake_post),
        patch(
            "worker.activities.document_generation.sources.download_to_tempfile",
            return_value=tmp_path,
        ),
    ):
        result = await hydrate_document_generation_sources(request)

    tmp_path.unlink(missing_ok=True)
    assert result.items[0].body == "Binary PDF source body"


@pytest.mark.asyncio
async def test_hydrate_sources_marks_image_only_pdf_as_scanned_fallback() -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "agent_file",
            "objectId": "00000000-0000-4000-8000-000000000031",
        }
    ]
    pdf = pymupdf.open()
    page = pdf.new_page(width=160, height=120)
    pixmap = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 80, 40), 0)
    pixmap.clear_with(220)
    page.insert_image(pymupdf.Rect(20, 20, 140, 90), pixmap=pixmap)
    with NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(pdf.tobytes())
        tmp_path = Path(tmp.name)
    pdf.close()

    try:
        async def fake_post(_path: str, body: dict) -> dict:
            return {
                "id": body["source"]["objectId"],
                "title": "Scanned PDF",
                "body": "scan.pdf (pdf, 4096 bytes)",
                "kind": "agent_file",
                "mimeType": "application/pdf",
                "objectKey": "agent-files/scan.pdf",
                "bytes": 4096,
            }

        with (
            patch("worker.activities.document_generation.sources.post_internal", fake_post),
            patch(
                "worker.activities.document_generation.sources.download_to_tempfile",
                return_value=tmp_path,
            ),
        ):
            result = await hydrate_document_generation_sources(request)

        assert result.items[0].body == "scan.pdf (pdf, 4096 bytes)"
        assert result.items[0].quality_signals == [
            "scanned_no_text",
            "metadata_fallback",
        ]
    finally:
        tmp_path.unlink(missing_ok=True)


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_keeps_metadata_fallback_for_unsupported_binary(
) -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "agent_file",
            "objectId": "00000000-0000-4000-8000-000000000031",
        }
    ]

    async def fake_post(_path: str, body: dict) -> dict:
        return {
            "id": body["source"]["objectId"],
            "title": "Opaque upload",
            "body": "archive.bin (binary, 4096 bytes)",
            "kind": "agent_file",
            "mimeType": "application/octet-stream",
            "objectKey": "agent-files/archive.bin",
            "bytes": 4096,
        }

    with (
        patch("worker.activities.document_generation.sources.post_internal", fake_post),
        patch("worker.activities.document_generation.sources.download_to_tempfile") as download,
    ):
        result = await hydrate_document_generation_sources(request)

    assert result.items[0].body == "archive.bin (binary, 4096 bytes)"
    download.assert_not_called()


@pytest.mark.asyncio
async def test_hydrate_document_generation_sources_skips_oversized_supported_binary() -> None:
    request = _request("pdf")
    request["generation"]["sources"] = [
        {
            "type": "agent_file",
            "objectId": "00000000-0000-4000-8000-000000000031",
        }
    ]

    async def fake_post(_path: str, body: dict) -> dict:
        return {
            "id": body["source"]["objectId"],
            "title": "Large PDF",
            "body": "large.pdf (pdf, 31457280 bytes)",
            "kind": "agent_file",
            "mimeType": "application/pdf",
            "objectKey": "agent-files/large.pdf",
            "bytes": 30 * 1024 * 1024,
        }

    with (
        patch("worker.activities.document_generation.sources.post_internal", fake_post),
        patch("worker.activities.document_generation.sources.download_to_tempfile") as download,
    ):
        result = await hydrate_document_generation_sources(request)

    assert result.items[0].body == "large.pdf (pdf, 31457280 bytes)"
    download.assert_not_called()


@pytest.mark.asyncio
async def test_generate_document_artifact_preserves_korean_and_long_pdf_content() -> None:
    uploaded: dict[str, object] = {}
    request = _request("pdf")
    request["generation"]["prompt"] = "\n".join(
        f"한글 문단 {index}: 긴 문서 생성 결과를 절단하지 않고 PDF에 보존합니다."
        for index in range(80)
    )

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(request)

    doc = pymupdf.open(stream=bytes(uploaded["data"]), filetype="pdf")
    extracted = "\n".join(page.get_text() for page in doc)
    assert len(doc) > 1
    assert "Objective" in extracted
    assert "한글 문단 0" in extracted
    assert "한글 문단 79" in extracted


@pytest.mark.asyncio
async def test_generate_document_artifact_pdf_preserves_long_unspaced_tokens() -> None:
    uploaded: dict[str, object] = {}
    request = _request("pdf")
    long_token = "긴토큰" * 600
    request["generation"]["prompt"] = long_token

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(request)

    doc = pymupdf.open(stream=bytes(uploaded["data"]), filetype="pdf")
    extracted = "".join(page.get_text() for page in doc).replace("\n", "").replace(" ", "")
    assert long_token in extracted


@pytest.mark.asyncio
async def test_generate_document_artifact_latex_pdf_posts_safe_tex_to_tectonic() -> None:
    uploaded: dict[str, object] = {}
    posted: dict[str, str] = {}
    request = _request("pdf")
    request["generation"]["template"] = "technical_report"
    request["generation"]["renderEngine"] = "latex"
    request["generation"]["prompt"] = "Generate a report with 50% growth & Korean text."

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    async def fake_post_tectonic(tex_source: str, bib_source: str) -> bytes:
        posted["tex"] = tex_source
        posted["bib"] = bib_source
        return b"%PDF-1.7\nlatex pdf"

    with (
        patch("worker.activities.document_generation.generate.upload_bytes", fake_upload),
        patch("worker.activities.document_generation.generate._post_tectonic", fake_post_tectonic),
    ):
        result = await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "source-a",
                        "title": "시장 & 조사_노트",
                        "body": "성장률은 50%이고 A_B 경로를 검토했습니다.",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    }
                ]
            },
        )

    assert result.objectKey.endswith("/project-report.pdf")
    assert result.mimeType == "application/pdf"
    assert bytes(uploaded["data"]).startswith(b"%PDF-1.7")
    assert "\\usepackage{kotex}" in posted["tex"]
    assert "\\tableofcontents" in posted["tex"]
    assert "\\section{Objective}" in posted["tex"]
    assert "50\\% growth \\& Korean text" in posted["tex"]
    assert "시장 \\& 조사\\_노트" in posted["tex"]
    assert "A\\_B" in posted["tex"]
    assert posted["bib"] == ""


@pytest.mark.asyncio
async def test_generate_document_artifact_creates_svg_figure() -> None:
    uploaded: dict[str, object] = {}
    request = _request("image")
    request["generation"]["template"] = "research_brief"
    request["generation"]["destination"]["filename"] = "evidence-map.svg"
    request["generation"]["prompt"] = "Map the evidence: 50% signal & source confidence."

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        result = await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "source-a",
                        "title": "시장 & 조사",
                        "body": "핵심 근거입니다.",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    }
                ]
            },
        )

    svg = bytes(uploaded["data"]).decode("utf-8")
    assert result.objectKey.endswith("/evidence-map.svg")
    assert result.mimeType == "image/svg+xml"
    assert uploaded["content_type"] == "image/svg+xml"
    assert svg.startswith("<svg")
    assert "Map the evidence: 50% signal &amp; source confidence." in svg
    assert "시장 &amp; 조사" in svg


@pytest.mark.asyncio
async def test_generate_document_artifact_model_image_uses_llm_provider() -> None:
    uploaded: dict[str, object] = {}
    request = _request("image")
    request["generation"]["imageEngine"] = "model"
    request["generation"]["destination"]["filename"] = "evidence-map.png"
    request["generation"]["prompt"] = "Make a clear strategy figure."

    class FakeProvider:
        async def generate_image(self, prompt: str):
            from types import SimpleNamespace

            assert "Make a clear strategy figure." in prompt
            assert "Sources:" in prompt
            return SimpleNamespace(
                image_bytes=b"\x89PNG\r\nmodel-image",
                mime_type="image/png",
                model="gemini-3.1-flash-image-preview",
            )

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with (
        patch("worker.activities.document_generation.generate.upload_bytes", fake_upload),
        patch("llm.factory.get_provider", return_value=FakeProvider()),
    ):
        result = await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "source-a",
                        "title": "시장 조사",
                        "body": "핵심 근거입니다.",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    }
                ]
            },
        )

    assert result.objectKey.endswith("/evidence-map.png")
    assert result.mimeType == "image/png"
    assert uploaded["content_type"] == "image/png"
    assert bytes(uploaded["data"]).startswith(b"\x89PNG")


@pytest.mark.asyncio
async def test_generate_document_artifact_creates_source_aware_docx() -> None:
    uploaded: dict[str, object] = {}
    request = _request("docx")

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        result = await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "00000000-0000-4000-8000-000000000021",
                        "title": "시장 조사 노트",
                        "body": "핵심 근거: 한국어 DOCX 본문입니다.",
                        "kind": "note",
                        "token_count": 16,
                        "included": True,
                    }
                ]
            },
        )

    assert result.objectKey.endswith("/project-report.docx")
    document = Document(BytesIO(bytes(uploaded["data"])))
    document_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    table_text = "\n".join(
        cell.text for table in document.tables for row in table.rows for cell in row.cells
    )
    assert "시장 조사 노트" in document_text
    assert "한국어 DOCX 본문입니다" in document_text
    assert "Template" in table_text
    assert "report" in table_text
    assert '<w:br w:type="page"/>' not in document.element.xml


@pytest.mark.asyncio
async def test_generate_document_artifact_docx_does_not_add_trailing_page_break() -> None:
    uploaded: dict[str, object] = {}
    request = _request("docx")

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "source-a",
                        "title": "첫 번째 소스",
                        "body": "첫 번째 본문",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    },
                    {
                        "id": "source-b",
                        "title": "두 번째 소스",
                        "body": "두 번째 본문",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    },
                ]
            },
        )

    document = Document(BytesIO(bytes(uploaded["data"])))
    document_xml = document.element.xml
    assert document_xml.count('<w:br w:type="page"/>') == 1
    trailing_break_xml = '<w:br w:type="page"/></w:r></w:p></w:body></w:document>'
    assert not document_xml.rstrip().endswith(trailing_break_xml)


@pytest.mark.asyncio
async def test_generate_document_artifact_docx_includes_source_quality_notes() -> None:
    uploaded: dict[str, object] = {}
    request = _request("docx")
    request["generation"]["template"] = "research_summary"

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "source-ok",
                        "title": "정상 소스",
                        "body": "검증된 근거 본문",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    },
                    {
                        "id": "source-fallback",
                        "title": "메타데이터 소스",
                        "body": "",
                        "kind": "agent_file",
                        "token_count": 50_000,
                        "included": False,
                        "quality_signals": ["source_token_budget_exceeded", "metadata_fallback"],
                    },
                ]
            },
        )

    document = Document(BytesIO(bytes(uploaded["data"])))
    document_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Evidence Quality Notes" in document_text
    assert "source_token_budget_exceeded" in document_text
    assert "Evidence Register" in document_text


@pytest.mark.asyncio
async def test_generate_document_artifact_creates_readable_pptx_deck() -> None:
    uploaded: dict[str, object] = {}

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(
            _request("pptx"),
            {
                "items": [
                    {
                        "id": "00000000-0000-4000-8000-000000000021",
                        "title": "시장 조사 노트",
                        "body": "첫 번째 근거\n두 번째 근거",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    }
                ]
            },
        )

    presentation = Presentation(BytesIO(bytes(uploaded["data"])))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert len(presentation.slides) >= 3
    assert "Project report" in slide_text
    assert "시장 조사 노트" in slide_text


@pytest.mark.asyncio
async def test_generate_document_artifact_pptx_preserves_long_sources() -> None:
    uploaded: dict[str, object] = {}
    sources = [
        {
            "id": f"00000000-0000-4000-8000-0000000001{index:02d}",
            "title": f"소스 {index}",
            "body": "\n".join(f"소스 {index} 라인 {line}" for line in range(10)),
            "kind": "note",
            "token_count": 20,
            "included": True,
        }
        for index in range(10)
    ]

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(_request("pptx"), {"items": sources})

    presentation = Presentation(BytesIO(bytes(uploaded["data"])))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert len(presentation.slides) < 30
    assert "소스 0 라인 9" in slide_text
    assert "소스 9 라인 9" in slide_text


@pytest.mark.asyncio
async def test_generate_document_artifact_creates_structured_xlsx_workbook() -> None:
    uploaded: dict[str, object] = {}

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(
            _request("xlsx"),
            {
                "items": [
                    {
                        "id": "00000000-0000-4000-8000-000000000021",
                        "title": "시장 조사 노트",
                        "body": "핵심 근거: 스프레드시트 검증",
                        "kind": "note",
                        "token_count": 10,
                        "included": True,
                    }
                ]
            },
        )

    workbook = load_workbook(BytesIO(bytes(uploaded["data"])))
    assert workbook.sheetnames == ["Summary", "Outline", "Sources"]
    summary = workbook["Summary"]
    outline = workbook["Outline"]
    sources = workbook["Sources"]
    assert summary["A1"].value == "Field"
    assert summary["B1"].value == "Value"
    assert outline["A1"].value == "Section"
    assert sources["A1"].value == "Source ID"
    assert sources["B2"].value == "note"
    assert sources["C2"].value == "시장 조사 노트"
    assert sources["D2"].value is True


@pytest.mark.asyncio
async def test_generate_document_artifact_xlsx_records_template_and_quality_signals() -> None:
    uploaded: dict[str, object] = {}
    request = _request("xlsx")
    request["generation"]["template"] = "spreadsheet"

    def fake_upload(key: str, data: bytes, content_type: str) -> str:
        uploaded.update({"key": key, "data": data, "content_type": content_type})
        return key

    with patch("worker.activities.document_generation.generate.upload_bytes", fake_upload):
        await generate_document_artifact(
            request,
            {
                "items": [
                    {
                        "id": "source-ok",
                        "title": "정상 데이터",
                        "body": "셀에 들어갈 본문",
                        "kind": "note",
                        "token_count": 8,
                        "included": True,
                    },
                    {
                        "id": "source-unsupported",
                        "title": "지원하지 않는 업로드",
                        "body": "archive.bin (binary, 4096 bytes)",
                        "kind": "agent_file",
                        "token_count": 12,
                        "included": True,
                        "quality_signals": ["unsupported_source", "metadata_fallback"],
                    },
                ]
            },
        )

    workbook = load_workbook(BytesIO(bytes(uploaded["data"])))
    assert workbook.sheetnames == ["Summary", "Outline", "Sources", "Quality"]
    summary_values = [row[1].value for row in workbook["Summary"].iter_rows(min_row=1, max_col=2)]
    assert "spreadsheet" in summary_values
    source_rows = list(workbook["Sources"].iter_rows(values_only=True))
    assert source_rows[0] == (
        "Source ID",
        "Kind",
        "Title",
        "Included",
        "Quality Signals",
        "Excerpt",
    )
    assert source_rows[2][4] == "unsupported_source, metadata_fallback"
    assert "unsupported_source" in workbook["Quality"]["A2"].value


@pytest.mark.asyncio
async def test_register_document_generation_result_posts_internal_agent_file_payload() -> None:
    calls: list[tuple[str, dict]] = []

    async def fake_post(path: str, body: dict) -> dict:
        calls.append((path, body))
        return {
            "object": {
                "id": "00000000-0000-4000-8000-000000000010",
                "objectType": "agent_file",
                "title": "Project report",
                "filename": "project-report.pdf",
                "kind": "pdf",
                "mimeType": "application/pdf",
                "projectId": "00000000-0000-4000-8000-000000000003",
            }
        }

    artifact = GeneratedDocumentArtifact(
        objectKey="agent-files/project/document-generation/request/project-report.pdf",
        mimeType="application/pdf",
        bytes=128,
    )
    with patch("worker.activities.document_generation.register.post_internal", fake_post):
        result = await register_document_generation_result(
            _request("pdf"),
            artifact,
            "document-generation/00000000-0000-4000-8000-000000000020",
        )

    assert result.id == "00000000-0000-4000-8000-000000000010"
    assert calls[0][0] == "/api/internal/document-generation/agent-files"
    body = calls[0][1]
    assert body["actionId"] == "00000000-0000-4000-8000-000000000011"
    assert body["requestId"] == "00000000-0000-4000-8000-000000000020"
    assert body["workflowId"] == "document-generation/00000000-0000-4000-8000-000000000020"
    assert body["workspaceId"] == "00000000-0000-4000-8000-000000000001"
    assert body["projectId"] == "00000000-0000-4000-8000-000000000003"
    assert body["objectKey"] == artifact.objectKey
    assert body["source"] == "document_generation"
